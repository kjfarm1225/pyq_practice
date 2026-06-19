# Streamlitダッシュボード：
# 工程内（取引先）・搬入・市場 MELMB LCM BUCK-IC 不具合発生状況
#
# 実行方法：
# 1) pip install streamlit pandas plotly openpyxl
# 2) streamlit run app.py

from io import BytesIO

import pandas as pd
import plotly.graph_objects as go
import streamlit as st
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt


# =========================
# 1. 画面設定J
# =========================
st.set_page_config(
    page_title="BUCK-IC 不具合発生状況",
    page_icon="📊",
    layout="wide",
)


# =========================
# 2. サンプルデータ作成
#    ※添付画像の表をもとに作成
# =========================
@st.cache_data
def load_sample_data() -> pd.DataFrame:
    data = [
        {
            "period": "23/2~23/3",
            "target": "対象1",
            "buck_ic_production": 209_644,
            "melmb_ng": 2,
            "melmb_ppm": 9.5,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 3,
            "incoming_ppm": 14.3,
            "market_ng": 8,
            "market_ppm": 38.2,
        },
        {
            "period": "23/3~23/4",
            "target": "対象2",
            "buck_ic_production": 213_767,
            "melmb_ng": 5,
            "melmb_ppm": 23.4,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 4,
            "incoming_ppm": 18.7,
            "market_ng": 9,
            "market_ppm": 42.1,
        },
        {
            "period": "23/4~23/5",
            "target": "対象3",
            "buck_ic_production": 270_160,
            "melmb_ng": 9,
            "melmb_ppm": 33.3,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 1,
            "incoming_ppm": 3.7,
            "market_ng": 13,
            "market_ppm": 48.1,
        },
        {
            "period": "23/5~23/7",
            "target": "対象4",
            "buck_ic_production": 695_655,
            "melmb_ng": 11,
            "melmb_ppm": 15.8,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 17,
            "incoming_ppm": 24.4,
            "market_ng": 21,
            "market_ppm": 30.2,
        },
        {
            "period": "23/7~23/9",
            "target": "対象5",
            "buck_ic_production": 634_397,
            "melmb_ng": 5,
            "melmb_ppm": 7.9,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 6,
            "incoming_ppm": 9.5,
            "market_ng": 11,
            "market_ppm": 17.3,
        },
        {
            "period": "23/9~24/2",
            "target": "TI対策1,2",
            "buck_ic_production": 2_073_018,
            "melmb_ng": 25,
            "melmb_ppm": 12.1,
            "sti_ng": 3,
            "sti_ppm": 1.4,
            "incoming_ng": 16,
            "incoming_ppm": 7.7,
            "market_ng": 60,
            "market_ppm": 28.9,
        },
        {
            "period": "24/2~24/4",
            "target": "TI対策3",
            "buck_ic_production": 643_470,
            "melmb_ng": 1,
            "melmb_ppm": 1.6,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 4,
            "incoming_ppm": 6.2,
            "market_ng": 9,
            "market_ppm": 14.0,
        },
        {
            "period": "24/4~24/6",
            "target": "対象7",
            "buck_ic_production": 1_007_135,
            "melmb_ng": 18,
            "melmb_ppm": 17.9,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 13,
            "incoming_ppm": 12.9,
            "market_ng": 12,
            "market_ppm": 11.9,
        },
        {
            "period": "24/6~24/7",
            "target": "対象8",
            "buck_ic_production": 469_036,
            "melmb_ng": 1,
            "melmb_ppm": 2.1,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 6,
            "incoming_ppm": 12.8,
            "market_ng": 4,
            "market_ppm": 8.5,
        },
        {
            "period": "24/7~24/10",
            "target": "TI対策4",
            "buck_ic_production": 1_003_576,
            "melmb_ng": 9,
            "melmb_ppm": 9.0,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 4,
            "incoming_ppm": 4.0,
            "market_ng": 11,
            "market_ppm": 11.0,
        },
        {
            "period": "24/11~24/12",
            "target": "TI対策5",
            "buck_ic_production": 246_709,
            "melmb_ng": 2,
            "melmb_ppm": 8.1,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 1,
            "incoming_ppm": 4.1,
            "market_ng": 3,
            "market_ppm": 12.2,
        },
        {
            "period": "24/12~25/1",
            "target": "TI対策5(M?)",
            "buck_ic_production": 751_382,
            "melmb_ng": 8,
            "melmb_ppm": 10.6,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 5,
            "incoming_ppm": 6.7,
            "market_ng": 15,
            "market_ppm": 20.0,
        },
        {
            "period": "25/1~26/3",
            "target": "TI対策6(設前)",
            "buck_ic_production": 3_578_399,
            "melmb_ng": 26,
            "melmb_ppm": 7.3,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 12,
            "incoming_ppm": 3.4,
            "market_ng": 18,
            "market_ppm": 5.0,
        },
        {
            "period": "25/1~26/3",
            "target": "LCM設変",
            "buck_ic_production": 1_565_723,
            "melmb_ng": 11,
            "melmb_ppm": 7.0,
            "sti_ng": 0,
            "sti_ppm": 0.0,
            "incoming_ng": 2,
            "incoming_ppm": 1.3,
            "market_ng": 2,
            "market_ppm": 1.3,
        },
    ]

    df = pd.DataFrame(data)
    df["x_label"] = df["target"] + "<br>" + df["period"]
    return df


# =========================
# 3. アップロードデータ読込
# =========================
def load_uploaded_file(uploaded_file) -> pd.DataFrame:
    if uploaded_file.name.endswith(".csv"):
        df = pd.read_csv(uploaded_file)
    else:
        df = pd.read_excel(uploaded_file)

    required_cols = [
        "period",
        "target",
        "buck_ic_production",
        "melmb_ng",
        "melmb_ppm",
        "sti_ng",
        "sti_ppm",
        "incoming_ng",
        "incoming_ppm",
        "market_ng",
        "market_ppm",
    ]

    missing_cols = [col for col in required_cols if col not in df.columns]
    if missing_cols:
        st.error("アップロードファイルに必要列がありません。")
        st.write("不足列:", missing_cols)
        st.stop()

    df = df.copy()
    df["x_label"] = df["target"].astype(str) + "<br>" + df["period"].astype(str)
    return df


# =========================
# 4. グラフ作成
# =========================
def make_line_chart(df: pd.DataFrame, selected_series: list[str]) -> go.Figure:
    series_map = {
        "MELMB工程内": {
            "ppm_col": "melmb_ppm",
            "color": "#2F80ED",
            "marker": "diamond",
        },
        "STI工程内": {
            "ppm_col": "sti_ppm",
            "color": "#EB5757",
            "marker": "square",
        },
        "搬入": {
            "ppm_col": "incoming_ppm",
            "color": "#8E8E8E",
            "marker": "triangle-up",
        },
        "市場": {
            "ppm_col": "market_ppm",
            "color": "#F2A900",
            "marker": "x",
        },
    }

    fig = go.Figure()

    for name in selected_series:
        ppm_col = series_map[name]["ppm_col"]

        fig.add_trace(
            go.Scatter(
                x=df["x_label"],
                y=df[ppm_col],
                mode="lines+markers+text",
                name=name,
                text=df[ppm_col].round(1),
                textposition="top center",
                line=dict(width=3, color=series_map[name]["color"]),
                marker=dict(
                    size=8,
                    symbol=series_map[name]["marker"],
                    color=series_map[name]["color"],
                ),
                hovertemplate=(
                    "<b>%{fullData.name}</b><br>"
                    "対象・期間：%{x}<br>"
                    "不良率：%{y:.1f} ppm"
                    "<extra></extra>"
                ),
            )
        )

    fig.update_layout(
        title=dict(
            text="BUCK-IC 不良発生状況",
            x=0.5,
            xanchor="center",
            font=dict(size=24),
        ),
        height=470,
        plot_bgcolor="#EAF4FF",
        paper_bgcolor="white",
        hovermode="x unified",
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.04,
            xanchor="center",
            x=0.5,
        ),
        margin=dict(l=30, r=30, t=90, b=40),
        xaxis=dict(
            tickfont=dict(size=11),
            showgrid=True,
            gridcolor="#D0D7DE",
        ),
        yaxis=dict(
            title="不良率（ppm）",
            range=[0, max(50, df[[c for c in df.columns if c.endswith("_ppm")]].max().max() * 1.15)],
            showgrid=True,
            gridcolor="#D0D7DE",
        ),
    )

    return fig


# =========================
# 5. 表示用テーブル作成
# =========================
def make_summary_table(df: pd.DataFrame) -> pd.DataFrame:
    table_rows = []

    table_rows.append(
        {
            "区分": "BUCK-IC生産数",
            **dict(zip(df["target"], df["buck_ic_production"])),
            "合計": df["buck_ic_production"].sum(),
        }
    )

    groups = [
        ("MELMB工程内", "melmb_ng", "melmb_ppm"),
        ("STI工程内", "sti_ng", "sti_ppm"),
        ("搬入", "incoming_ng", "incoming_ppm"),
        ("市場", "market_ng", "market_ppm"),
    ]

    for group_name, ng_col, ppm_col in groups:
        table_rows.append(
            {
                "区分": f"{group_name}：不良数",
                **dict(zip(df["target"], df[ng_col])),
                "合計": df[ng_col].sum(),
            }
        )

        total_ng = df[ng_col].sum()
        total_production = df["buck_ic_production"].sum()
        total_ppm = total_ng / total_production * 1_000_000 if total_production else 0

        table_rows.append(
            {
                "区分": f"{group_name}：不良率(ppm)",
                **dict(zip(df["target"], df[ppm_col])),
                "合計": round(total_ppm, 1),
            }
        )

    summary = pd.DataFrame(table_rows)
    return summary


def style_summary_table(summary: pd.DataFrame):
    def color_rows(row):
        if "不良率" in str(row["区分"]):
            return ["background-color: #FFF7CC; color: blue; font-weight: bold"] * len(row)
        if "不良数" in str(row["区分"]):
            return ["background-color: #FFFBE6"] * len(row)
        if "BUCK-IC" in str(row["区分"]):
            return ["background-color: #F2F2F2; font-weight: bold"] * len(row)
        return [""] * len(row)

    return (
        summary.style
        .apply(color_rows, axis=1)
        .format(precision=1, thousands=",")
    )


def format_ppt_value(value) -> str:
    if pd.isna(value):
        return ""
    if isinstance(value, float) and not value.is_integer():
        return f"{value:,.1f}"
    if isinstance(value, float):
        return f"{int(value):,}"
    if isinstance(value, int):
        return f"{value:,}"
    return str(value)


def build_powerpoint(
    df: pd.DataFrame,
    summary: pd.DataFrame,
    data_note: str,
    selected_series: list[str],
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "工程内（取引先）・搬入・市場 不具合発生状況"
    title_slide.placeholders[1].text = (
        f"{data_note}\n"
        f"表示系列: {', '.join(selected_series)}\n"
        f"更新日: 2026年01月14日"
    )

    kpi_slide = prs.slides.add_slide(prs.slide_layouts[5])
    kpi_slide.shapes.title.text = "KPI"

    metrics = [
        ("BUCK-IC生産数 合計", df["buck_ic_production"].sum()),
        ("MELMB工程内 不良数", df["melmb_ng"].sum()),
        ("STI工程内 不良数", df["sti_ng"].sum()),
        ("搬入 不良数", df["incoming_ng"].sum()),
        ("市場 不良数", df["market_ng"].sum()),
    ]

    rows = 2
    cols = len(metrics)
    left = Inches(0.45)
    top = Inches(1.6)
    width = Inches(12.4)
    height = Inches(1.8)
    table = kpi_slide.shapes.add_table(rows, cols, left, top, width, height).table

    for col_idx, (label, value) in enumerate(metrics):
        header_cell = table.cell(0, col_idx)
        value_cell = table.cell(1, col_idx)
        header_cell.text = label
        value_cell.text = format_ppt_value(value)
        for cell in (header_cell, value_cell):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
        for paragraph in header_cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True

    summary_slide = prs.slides.add_slide(prs.slide_layouts[5])
    summary_slide.shapes.title.text = "不良率(ppm) 集計表"

    summary_rows, summary_cols = summary.shape[0] + 1, summary.shape[1]
    table = summary_slide.shapes.add_table(
        summary_rows,
        summary_cols,
        Inches(0.25),
        Inches(1.35),
        Inches(12.85),
        Inches(5.85),
    ).table

    column_widths = [Inches(1.85)] + [Inches(0.85)] * (summary_cols - 2) + [Inches(1.0)]
    for idx, width_value in enumerate(column_widths):
        table.columns[idx].width = width_value

    header_fill = (242, 242, 242)
    row_colors = {
        "BUCK-IC": (242, 242, 242),
        "不良数": (255, 251, 230),
        "不良率": (255, 247, 204),
    }

    for col_idx, header in enumerate(summary.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*header_fill)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(9)

    for row_idx, (_, row) in enumerate(summary.iterrows(), start=1):
        row_fill = None
        row_label = str(row["区分"])
        for key, color in row_colors.items():
            if key in row_label:
                row_fill = color
                break

        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = format_ppt_value(value)
            if row_fill is not None:
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*row_fill)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)
                    if "不良率" in row_label or "BUCK-IC" in row_label:
                        run.font.bold = True

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


# =========================
# 6. サイドバー
# =========================
st.sidebar.header("表示設定")

uploaded_file = st.sidebar.file_uploader(
    "CSV / Excelをアップロード",
    type=["csv", "xlsx"],
    help="列名がサンプルデータと同じ形式であれば、手元データに差し替えできます。",
)

if uploaded_file is None:
    df = load_sample_data()
    data_note = "サンプルデータ表示中（添付画像ベース）"
else:
    df = load_uploaded_file(uploaded_file)
    data_note = f"アップロードデータ表示中：{uploaded_file.name}"

all_series = ["MELMB工程内", "STI工程内", "搬入", "市場"]
selected_series = st.sidebar.multiselect(
    "グラフに表示する系列",
    options=all_series,
    default=all_series,
)

show_table = st.sidebar.checkbox("詳細テーブルを表示", value=True)
show_raw_data = st.sidebar.checkbox("元データを表示", value=False)


# =========================
# 7. ヘッダー
# =========================
left, right = st.columns([3, 1])

with left:
    st.title("工程内（取引先）・搬入・市場　MELMB LCM BUCK-IC 不具合発生状況")
    st.caption(data_note)

with right:
    st.markdown(
        """
        <div style="text-align:right; line-height:1.7;">
            <b>更新日：</b>2026年01月14日<br>
            <b>作成日：</b>2024年08月02日<br>
            品質本部　調達品質部　梶間
        </div>
        """,
        unsafe_allow_html=True,
    )


# =========================
# 8. KPI
# =========================
total_production = df["buck_ic_production"].sum()
total_melmb_ng = df["melmb_ng"].sum()
total_sti_ng = df["sti_ng"].sum()
total_incoming_ng = df["incoming_ng"].sum()
total_market_ng = df["market_ng"].sum()

kpi1, kpi2, kpi3, kpi4, kpi5 = st.columns(5)

kpi1.metric("BUCK-IC生産数 合計", f"{total_production:,}")
kpi2.metric("MELMB工程内 不良数", f"{total_melmb_ng:,}")
kpi3.metric("STI工程内 不良数", f"{total_sti_ng:,}")
kpi4.metric("搬入 不良数", f"{total_incoming_ng:,}")
kpi5.metric("市場 不良数", f"{total_market_ng:,}")


# =========================
# 9. グラフ
# =========================
st.plotly_chart(
    make_line_chart(df, selected_series),
    use_container_width=True,
)


# =========================
# 10. 詳細テーブル
# =========================
if show_table:
    st.subheader("不良率(ppm) ＝ 不良数 / BUCK-IC生産数")

    summary = make_summary_table(df)

    st.dataframe(
        style_summary_table(summary),
        use_container_width=True,
        hide_index=True,
        height=390,
    )

    csv = summary.to_csv(index=False).encode("utf-8-sig")
    pptx = build_powerpoint(df, summary, data_note, selected_series)
    csv_col, ppt_col = st.columns(2)

    with csv_col:
        st.download_button(
            label="集計表をCSVでダウンロード",
            data=csv,
            file_name="buck_ic_summary.csv",
            mime="text/csv",
        )

    with ppt_col:
        st.download_button(
            label="PowerPointでダウンロード",
            data=pptx,
            file_name="buck_ic_summary.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )


# =========================
# 11. 元データ
# =========================
if show_raw_data:
    st.subheader("元データ")
    st.dataframe(df, use_container_width=True, hide_index=True)
